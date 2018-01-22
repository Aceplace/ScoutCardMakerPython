from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Pt
import csv



PLAYER_WIDTH = 265265
PLAYER_HEIGHT = 220540
LINEMAN_TOP = 5467068
LT_LEFT = 3761549
LG_LEFT = 4129594
C_LEFT = 4464179
RG_LEFT = 4810890
RT_LEFT = 5178935

LEFT_HASH_OFFSET = -1615739
RIGHT_HASH_OFFSET = 1553599

FIELD_URL = "field1.png"
FIELD_LEFT = 0
FIELD_TOP = 716946
FIELD_WIDTH = 9144000
FIELD_HEIGHT = 4766654

TITLE_LEFT = 106531
TITLE_TOP = 132963
TITLE_WIDTH = 8904303
TITLE_HEIGHT = 369332


#plays follow convention
#play[0]->play
#play[1]->hash
#play[2]->formation
#play[3]->backfield
#play[4]->play
def formation_card_maker(script_file, formation_templates_file, backfield_templates_file):

    plays = None
    with open(script_file) as csvfile:
        csvreader = csv.reader(csvfile)
        csvreader.__next__()
        plays = [row for row in csvreader]

    formation_templates_prs = Presentation(formation_templates_file)
    backfield_templates_prs = Presentation(backfield_templates_file)
    prs = Presentation()

    for play in plays:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture("field1.png", 0, 716946, 9144000, 4766654)
        add_lineman(slide, play[1])
        add_backfield(slide, play[1], play[3], backfield_templates_prs)
        add_formation(slide, play[1], play[2], formation_templates_prs)
        text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
        text_box.text_frame.text = play[0] + ' ' + play[1] + ' ' + play[2]+ ' ' + play[3]+ ' ' + play[4]
        text_box.text_frame.paragraphs[0].font.size = Pt(24)

    prs.save('test.pptx')



def add_lineman(prs_slide, hash_mark):
    offset = 0
    if hash_mark == 'l' or hash_mark == 'L':
        offset = LEFT_HASH_OFFSET
    if hash_mark == 'r' or hash_mark == 'R':
        offset = RIGHT_HASH_OFFSET
    shape = prs_slide.shapes.add_shape(MSO_SHAPE.OVAL, LT_LEFT + offset, LINEMAN_TOP, PLAYER_WIDTH, PLAYER_HEIGHT)
    lineman_fill(shape)
    shape = prs_slide.shapes.add_shape(MSO_SHAPE.OVAL, LG_LEFT + offset, LINEMAN_TOP, PLAYER_WIDTH, PLAYER_HEIGHT)
    lineman_fill(shape)
    shape = prs_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, C_LEFT + offset, LINEMAN_TOP, PLAYER_WIDTH, PLAYER_HEIGHT)
    lineman_fill(shape)
    shape = prs_slide.shapes.add_shape(MSO_SHAPE.OVAL, RG_LEFT + offset, LINEMAN_TOP, PLAYER_WIDTH, PLAYER_HEIGHT)
    lineman_fill(shape)
    shape = prs_slide.shapes.add_shape(MSO_SHAPE.OVAL, RT_LEFT + offset, LINEMAN_TOP, PLAYER_WIDTH, PLAYER_HEIGHT)
    lineman_fill(shape)


def add_backfield(prs_slide, hash_mark, back_field, backfield_templates_prs):
    offset = 0
    if hash_mark == 'l' or hash_mark == 'L':
        offset = LEFT_HASH_OFFSET
    if hash_mark == 'r' or hash_mark == 'R':
        offset = RIGHT_HASH_OFFSET

    #look for matching backfield in backfield_templates_prs
    matching_slide = None
    for slide in backfield_templates_prs.slides:
        text_boxes_texts = [shape.text_frame.text for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
        for text in text_boxes_texts:
            split_text = text.split('=')
            if len(split_text) > 1 and split_text[0].strip().upper() == 'NAME' and split_text[1].strip().upper() == back_field.strip().upper():
                matching_slide = slide
        if matching_slide != None:
            break

    #found matching slide, get position of skill players and add them to slide
    if matching_slide != None:
        skill_players_shapes = [shape for shape in matching_slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.text_frame.text != '']
        for skill_player_shape in skill_players_shapes:
            left = skill_player_shape.left
            top = skill_player_shape.top
            shape = prs_slide.shapes.add_shape(MSO_SHAPE.OVAL, left + offset, top, PLAYER_WIDTH, PLAYER_HEIGHT)
            skill_fill(shape, skill_player_shape.text_frame.text)


def add_formation(prs_slide, hash_mark, formation, formation_templates_prs):
    #look for matching backfield in backfield_templates_prs
    matching_slide = None
    found_formation = False
    found_hash = False
    for slide in formation_templates_prs.slides:
        text_boxes_texts = [shape.text_frame.text for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
        for text in text_boxes_texts:
            split_text = text.split('=')
            if len(split_text) > 1 and split_text[0].strip().upper() == 'NAME' and split_text[1].strip().upper() == formation.strip().upper():
                found_formation = True
            if len(split_text) > 1 and split_text[0].strip().upper() == 'HASH' and split_text[1].strip().upper() == hash_mark.strip().upper():
                found_hash = True
        if found_formation == True and found_hash == True:
            matching_slide = slide
            break
        found_formation = found_hash = False

    #found matching slide, get position of skill players and add them to slide
    if matching_slide != None:
        skill_players_shapes = [shape for shape in matching_slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.text_frame.text != '']
        for skill_player_shape in skill_players_shapes:
            left = skill_player_shape.left
            top = skill_player_shape.top
            shape = prs_slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, PLAYER_WIDTH, PLAYER_HEIGHT)
            skill_fill(shape, skill_player_shape.text_frame.text)


def lineman_fill(shape):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(50, 50, 156)
    shape.line.color.rgb = RGBColor(50, 50, 156)
    shape.line.color.brightness = -0.5
    shape.line.width = Pt(1.0)

def skill_fill(shape, label):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.width = Pt(1.0)
    shape.text_frame.text = label
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)




if __name__ == "__main__":
    formation_card_maker("Book1.csv", "formation_templates.pptx", "backfield_templates.pptx")
