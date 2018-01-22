from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.shapes.connector import Connector
from pptx.util import Pt
import csv


PLAYER_WIDTH = 633010
PLAYER_HEIGHT = 407504
LINEMAN_TOP = 4236278
LT_LEFT = 2492072
LG_LEFT = 3371985
C_LEFT = 4248754
RG_LEFT = 5125522
RT_LEFT = 6002291
CARD_CENTER = C_LEFT + PLAYER_WIDTH / 2
CONNECTOR_CENTER = 4572000

TITLE_LEFT = 106531
TITLE_TOP = 132963
TITLE_WIDTH = 8904303
TITLE_HEIGHT = 369332

DEFENSIVE_PLAYER_SIZE = Pt(40)
DEFENSIVE_PLAYER_FONT_NAME = 'Calibri'

#plays csv files follow convention
#play[0]->play
#play[1]->hash
#play[2]->formation
#play[3]->backfield
#play[4]->play
#play[5]->Scheme
#play[6]->Def front
#play[7]->Defensive Play Call

#formation core csv files follow convention
#core_formations[0]->Formation Name
#core_formations[1]->Core_Name
def core_card_maker(script_file, core_templates_file, formation_core_file):

    plays = None
    with open(script_file) as csvfile:
        csvreader = csv.reader(csvfile)
        csvreader.__next__()
        plays = [row for row in csvreader]

    core_formations = None
    with open(formation_core_file) as csvfile:
        csvreader = csv.reader(csvfile)
        csvreader.__next__()
        core_formations = [row for row in csvreader]

    core_templates_prs = Presentation(core_templates_file)
    prs = Presentation()

    for play in plays:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_lineman(slide)
        core = get_core_from_formation(play[2], core_formations)
        add_core(slide, core, play[5], play[6], core_templates_prs)
        text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
        text_box.text_frame.text = play[0] + ' ' + play[1] + ' ' + play[2]+ ' ' + play[3]+ ' ' + play[4]
        text_box.text_frame.paragraphs[0].font.size = Pt(24)

    prs.save('test.pptx')



def add_lineman(prs_slide):
    shape = prs_slide.shapes.add_shape(MSO_SHAPE.OVAL, LT_LEFT, LINEMAN_TOP, PLAYER_WIDTH, PLAYER_HEIGHT)
    lineman_fill(shape)
    shape = prs_slide.shapes.add_shape(MSO_SHAPE.OVAL, LG_LEFT, LINEMAN_TOP, PLAYER_WIDTH, PLAYER_HEIGHT)
    lineman_fill(shape)
    shape = prs_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, C_LEFT, LINEMAN_TOP, PLAYER_WIDTH, PLAYER_HEIGHT)
    lineman_fill(shape)
    shape = prs_slide.shapes.add_shape(MSO_SHAPE.OVAL, RG_LEFT, LINEMAN_TOP, PLAYER_WIDTH, PLAYER_HEIGHT)
    lineman_fill(shape)
    shape = prs_slide.shapes.add_shape(MSO_SHAPE.OVAL, RT_LEFT, LINEMAN_TOP, PLAYER_WIDTH, PLAYER_HEIGHT)
    lineman_fill(shape)

def get_core_from_formation(formation_name, core_formations):
    for core_formation in core_formations:
        if formation_name.strip().upper() == core_formation[0].strip().upper():
            return core_formation[1]
    return None

def add_core(prs_slide, core, scheme, front, core_templates_prs):
    if core==None:
        return
    #look for matching core in core_templates_prs
    matching_slide = None
    found_core = False
    found_scheme = False
    found_front = False
    #flip the play if the last word in name is "LT"
    flip = core.strip()[-2:].upper() == 'LT'
    if flip:
        core = core.strip()[0:-2] + 'RT'
    for slide in core_templates_prs.slides:
        text_boxes_texts = [shape.text_frame.text for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
        for text in text_boxes_texts:
            split_text = text.split('=')
            if len(split_text) > 1 and split_text[0].strip().upper() == 'CORE' and split_text[1].strip().upper() == core.strip().upper():
                found_core = True
            if len(split_text) > 1 and split_text[0].strip().upper() == 'SCHEME' and split_text[1].strip().upper() == scheme.strip().upper():
                found_scheme = True
            if len(split_text) > 1 and split_text[0].strip().upper() == 'FRONT' and split_text[1].strip().upper() == front.strip().upper():
                found_front = True
        if found_core == True and found_scheme == True and found_front ==True:
            matching_slide = slide
            break
        found_core = found_scheme = found_front = False

    #found matching slide,
    if matching_slide != None:
        #get position of skill players and add them to slide
        skill_players_shapes = [shape for shape in matching_slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.text_frame.text != '']
        for skill_player_shape in skill_players_shapes:
            left = flip_position(skill_player_shape.left, PLAYER_WIDTH, flip)
            top = skill_player_shape.top
            shape = prs_slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, PLAYER_WIDTH, PLAYER_HEIGHT)
            skill_fill(shape, skill_player_shape.text_frame.text)

        #get defensive player text boxes (defensive players and add them to the slide)
        def_players_text_boxes = [shape for shape in matching_slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and len(shape.text_frame.text) == 1]
        for def_player_text_box in def_players_text_boxes:
            width = def_player_text_box.width
            height = def_player_text_box.height
            left = flip_position(def_player_text_box.left, width, flip)
            top = def_player_text_box.top
            shape = prs_slide.shapes.add_textbox(left, top, width, height)
            shape.text_frame.text = def_player_text_box.text_frame.text
            shape.text_frame.paragraphs[0].font.size = DEFENSIVE_PLAYER_SIZE
            shape.text_frame.paragraphs[0].font.name = DEFENSIVE_PLAYER_FONT_NAME

        #connectors
        #line_connects = [shape for shape in matching_slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.LINE]
        connectors = [shape for shape in matching_slide.shapes if isinstance(shape, Connector)]
        for connector in connectors:
            begin_x = max(0, flip_connector(connector.begin_x, flip))
            end_x = max(0, flip_connector(connector.end_x, flip))
            print("flip " + str(flip))
            print(begin_x, end_x)
            prs_slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, begin_x, connector.begin_y, end_x, connector.end_y)



def flip_position(position, width,  flip):
    if flip:
        if position >= CARD_CENTER:#to the right of center
            return CARD_CENTER - (position - CARD_CENTER) - width
        else: #to the left of center
            return CARD_CENTER + (CARD_CENTER - position) - width
    return position

def flip_connector(position, flip):
    if flip:
        if position >= CONNECTOR_CENTER:#to the right of center
            return CONNECTOR_CENTER - (position - CONNECTOR_CENTER)
        else: #to the left of center
            return CONNECTOR_CENTER + (CONNECTOR_CENTER - position)
    return position

def lineman_fill(shape):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.width = Pt(1.0)

def skill_fill(shape, label):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.line.width = Pt(1.0)
    shape.text_frame.text = label
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)




if __name__ == "__main__":
    core_card_maker("Book1.csv", "core_templates.pptx", "formation_to_core.csv")

    #prs = Presentation("core_templates.pptx")
    #slide = prs.slides[2]
    #for shape in slide.shapes:
    #    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
    #        print(shape.shadow)
