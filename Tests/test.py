from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches, Pt

prs = Presentation('test.pptx')
slide = prs.slides.add_slide(prs.slide_layouts[6])

slide.shapes.add_picture("field1.png", 0, 716946, 9144000, 4766654)


prs.save('test.pptx')
