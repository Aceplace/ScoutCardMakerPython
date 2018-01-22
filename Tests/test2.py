from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation("templateA.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            dimensions = shape.left, shape.top, shape.width, shape.height
            print(shape.shape_type)
            print(dimensions)
